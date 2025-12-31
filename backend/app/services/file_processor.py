"""
File processing service for extracting text from various file formats.
"""
import os
from pathlib import Path
from typing import Optional, List
import logging

logger = logging.getLogger(__name__)


class FileProcessor:
    """Service for processing uploaded files."""
    
    @staticmethod
    def get_file_type(filename: str) -> str:
        """Determine file type from filename."""
        ext = Path(filename).suffix.lower()

        type_mapping = {
            '.pdf': 'pdf',
            '.doc': 'doc',
            '.docx': 'docx',
            '.ppt': 'pptx',
            '.pptx': 'pptx',
            '.txt': 'text',
            '.md': 'markdown',
            '.png': 'image',
            '.jpg': 'image',
            '.jpeg': 'image',
            '.gif': 'image',
            '.webp': 'image',
        }

        return type_mapping.get(ext, 'unknown')
    
    @staticmethod
    async def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            from pypdf import PdfReader
            
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n\n"
            
            logger.info(f"Extracted {len(text)} characters from PDF: {file_path}")
            return text.strip()
        
        except Exception as e:
            logger.error(f"Error extracting text from PDF {file_path}: {str(e)}")
            raise
    
    @staticmethod
    async def extract_text_from_docx(file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            from docx import Document

            doc = Document(file_path)
            text = "\n\n".join([paragraph.text for paragraph in doc.paragraphs])

            logger.info(f"Extracted {len(text)} characters from DOCX: {file_path}")
            return text.strip()

        except Exception as e:
            logger.error(f"Error extracting text from DOCX {file_path}: {str(e)}")
            raise

    @staticmethod
    async def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PPTX file."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_runs = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)

            text = "\n\n".join(text_runs)
            logger.info(f"Extracted {len(text)} characters from PPTX: {file_path}")
            return text.strip()

        except Exception as e:
            logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
            raise

    @staticmethod
    async def extract_images_from_pptx(file_path: str) -> List[str]:
        """
        Extract images from PPTX file and save them to temporary directory.

        Returns:
            List of image file paths
        """
        try:
            from pptx import Presentation
            import tempfile
            from pathlib import Path

            prs = Presentation(file_path)
            image_paths = []

            # Create temp directory for extracted images
            temp_dir = tempfile.mkdtemp(prefix="pptx_images_")

            image_count = 0
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    # Check if shape has image
                    if hasattr(shape, "image"):
                        try:
                            image = shape.image
                            # Get image bytes
                            image_bytes = image.blob
                            # Get image extension
                            ext = image.ext

                            # Save image
                            image_filename = f"slide_{slide_idx + 1}_img_{image_count + 1}.{ext}"
                            image_path = os.path.join(temp_dir, image_filename)

                            with open(image_path, 'wb') as f:
                                f.write(image_bytes)

                            image_paths.append(image_path)
                            image_count += 1
                            logger.info(f"Extracted image: {image_filename}")
                        except Exception as img_error:
                            logger.warning(f"Failed to extract image from slide {slide_idx + 1}: {img_error}")
                            continue

            logger.info(f"Extracted {len(image_paths)} images from PPTX: {file_path}")
            return image_paths

        except Exception as e:
            logger.error(f"Error extracting images from PPTX {file_path}: {str(e)}")
            return []

    @staticmethod
    async def extract_text_from_text_file(file_path: str) -> str:
        """Extract text from plain text or markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            logger.info(f"Read {len(text)} characters from text file: {file_path}")
            return text.strip()
        
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {str(e)}")
            raise
    
    @staticmethod
    async def process_file(file_path: str, file_type: str) -> Optional[str]:
        """
        Process file and extract text content.
        
        Args:
            file_path: Path to the file
            file_type: Type of file (pdf, docx, text, markdown, image)
        
        Returns:
            Extracted text content, or None for images (will be processed by Gemini directly)
        """
        try:
            if file_type == 'pdf':
                return await FileProcessor.extract_text_from_pdf(file_path)

            elif file_type == 'docx':
                return await FileProcessor.extract_text_from_docx(file_path)

            elif file_type == 'pptx':
                return await FileProcessor.extract_text_from_pptx(file_path)

            elif file_type in ['text', 'markdown']:
                return await FileProcessor.extract_text_from_text_file(file_path)

            elif file_type == 'image':
                # Images will be processed directly by Gemini's vision API
                logger.info(f"Image file will be processed by Gemini: {file_path}")
                return None
            
            else:
                logger.warning(f"Unsupported file type: {file_type}")
                return None
        
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            raise


# Global instance
file_processor = FileProcessor()

